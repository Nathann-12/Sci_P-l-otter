"""Publication report model + renderers (HTML / Markdown / PDF).

Pure and Qt-free: the UI/AI layers assemble a :class:`ReportDocument` from open
Graphs, result Books and analysis provenance, and this module renders it. The
HTML output is fully self-contained (images embedded as base64) so a single
file opens anywhere and prints to a clean PDF; Markdown is for version control;
PDF is a direct reportlab render.
"""
from __future__ import annotations

import base64
import datetime as _dt
import html
import io
from dataclasses import dataclass, field
from typing import List, Optional, Sequence

import pandas as pd


# --------------------------------------------------------------- document model
@dataclass
class TextSection:
    heading: str = ""
    body: str = ""
    kind: str = "text"


@dataclass
class FigureSection:
    title: str = ""
    png: bytes = b""
    caption: str = ""
    kind: str = "figure"


@dataclass
class TableSection:
    title: str = ""
    frame: Optional[pd.DataFrame] = None
    caption: str = ""
    max_rows: int = 30
    kind: str = "table"


@dataclass
class KpiSection:
    items: List[tuple] = field(default_factory=list)  # (label, value)
    kind: str = "kpi"


TEMPLATES = ("Lab Report", "Journal Manuscript", "Thesis Chapter", "Presentation Handout")


@dataclass
class ReportDocument:
    title: str = "SciPlotter Report"
    subtitle: str = ""
    author: str = ""
    date: str = ""
    template: str = "Lab Report"
    sections: List[object] = field(default_factory=list)

    def __post_init__(self):
        if not self.date:
            self.date = _dt.date.today().isoformat()

    # small builder helpers so callers read cleanly
    def add_text(self, heading: str, body: str) -> "ReportDocument":
        self.sections.append(TextSection(heading, body))
        return self

    def add_figure(self, title: str, png: bytes, caption: str = "") -> "ReportDocument":
        self.sections.append(FigureSection(title, png, caption))
        return self

    def add_table(self, title: str, frame: pd.DataFrame, caption: str = "",
                  max_rows: int = 30) -> "ReportDocument":
        self.sections.append(TableSection(title, frame, caption, max_rows))
        return self

    def add_kpis(self, items: Sequence[tuple]) -> "ReportDocument":
        self.sections.append(KpiSection(list(items)))
        return self

    def counts(self) -> dict:
        return {
            "figures": sum(isinstance(s, FigureSection) for s in self.sections),
            "tables": sum(isinstance(s, TableSection) for s in self.sections),
            "text": sum(isinstance(s, TextSection) for s in self.sections),
        }


# --------------------------------------------------------------------- helpers
def _fmt(value) -> str:
    if isinstance(value, float):
        if value != value:  # NaN
            return "—"
        if value == 0 or (1e-4 <= abs(value) < 1e6):
            return f"{value:.4g}"
        return f"{value:.4e}"
    return str(value)


def _table_html(frame: pd.DataFrame, max_rows: int) -> str:
    shown = frame.head(max_rows)
    head = "".join(f"<th>{html.escape(str(c))}</th>" for c in shown.columns)
    rows = []
    for _, row in shown.iterrows():
        cells = "".join(f"<td>{html.escape(_fmt(v))}</td>" for v in row)
        rows.append(f"<tr>{cells}</tr>")
    more = ""
    if len(frame) > max_rows:
        more = (f'<p class="more">… {len(frame) - max_rows} more rows '
                f'({len(frame)} total)</p>')
    return (f'<table><thead><tr>{head}</tr></thead>'
            f'<tbody>{"".join(rows)}</tbody></table>{more}')


# ------------------------------------------------------------------------- HTML
_CSS = """
:root{--ink:#1a2230;--muted:#5a6676;--line:#e2e7ee;--accent:#2f7fe0;--bg:#fff;--soft:#f6f8fb;}
*{box-sizing:border-box;}
body{margin:0;background:var(--soft);color:var(--ink);
 font-family:"Segoe UI","Sarabun","Leelawadee UI",system-ui,-apple-system,"Noto Sans Thai",sans-serif;
 line-height:1.6;}
.page{max-width:820px;margin:28px auto;background:var(--bg);padding:56px 64px;
 box-shadow:0 1px 3px rgba(20,30,50,.08),0 12px 32px rgba(20,30,50,.07);border-radius:6px;}
header.doc{border-bottom:2px solid var(--accent);padding-bottom:18px;margin-bottom:28px;}
h1.title{font-size:30px;line-height:1.2;margin:0 0 6px;letter-spacing:-.01em;}
.subtitle{color:var(--muted);font-size:17px;margin:0 0 14px;}
.meta{display:flex;gap:22px;flex-wrap:wrap;color:var(--muted);font-size:13px;}
.meta b{color:var(--ink);font-weight:600;}
section{margin:26px 0;}
h2{font-size:19px;margin:0 0 10px;padding-bottom:4px;border-bottom:1px solid var(--line);}
p{margin:0 0 12px;}
figure{margin:0;text-align:center;}
figure img{max-width:100%;height:auto;border:1px solid var(--line);border-radius:4px;}
figcaption{color:var(--muted);font-size:13px;margin-top:8px;text-align:center;font-style:italic;}
.tablewrap{overflow-x:auto;}
table{border-collapse:collapse;width:100%;font-size:13px;font-variant-numeric:tabular-nums;}
th,td{padding:7px 11px;border-bottom:1px solid var(--line);text-align:left;}
th{background:var(--soft);font-weight:600;border-bottom:2px solid var(--line);}
tbody tr:hover{background:var(--soft);}
.more{color:var(--muted);font-size:12px;margin-top:6px;}
.caption{color:var(--muted);font-size:13px;font-style:italic;margin-top:6px;}
.kpis{display:grid;grid-template-columns:repeat(auto-fit,minmax(150px,1fr));gap:12px;}
.kpi{background:var(--soft);border:1px solid var(--line);border-radius:8px;padding:14px 16px;}
.kpi .v{font-size:22px;font-weight:700;letter-spacing:-.01em;}
.kpi .l{color:var(--muted);font-size:12px;margin-top:4px;}
footer.doc{margin-top:34px;padding-top:14px;border-top:1px solid var(--line);
 color:var(--muted);font-size:12px;}
@media print{body{background:#fff;}.page{box-shadow:none;margin:0;max-width:none;
 padding:0;border-radius:0;}@page{margin:1.8cm;}}
"""


def render_html(doc: ReportDocument) -> str:
    """Return a single self-contained HTML string (images embedded)."""
    parts: List[str] = []
    for s in doc.sections:
        if isinstance(s, TextSection):
            block = []
            if s.heading:
                block.append(f"<h2>{html.escape(s.heading)}</h2>")
            for para in str(s.body).split("\n\n"):
                if para.strip():
                    block.append(f"<p>{html.escape(para.strip())}</p>")
            parts.append(f"<section>{''.join(block)}</section>")
        elif isinstance(s, FigureSection):
            b64 = base64.b64encode(s.png).decode("ascii") if s.png else ""
            img = (f'<img src="data:image/png;base64,{b64}" alt="{html.escape(s.title)}">'
                   if b64 else "")
            cap = f"<figcaption>{html.escape(s.caption)}</figcaption>" if s.caption else ""
            title = f"<h2>{html.escape(s.title)}</h2>" if s.title else ""
            parts.append(f"<section>{title}<figure>{img}{cap}</figure></section>")
        elif isinstance(s, TableSection) and s.frame is not None:
            title = f"<h2>{html.escape(s.title)}</h2>" if s.title else ""
            cap = f'<p class="caption">{html.escape(s.caption)}</p>' if s.caption else ""
            parts.append(
                f'<section>{title}<div class="tablewrap">'
                f'{_table_html(s.frame, s.max_rows)}</div>{cap}</section>')
        elif isinstance(s, KpiSection) and s.items:
            cards = "".join(
                f'<div class="kpi"><div class="v">{html.escape(_fmt(v))}</div>'
                f'<div class="l">{html.escape(str(l))}</div></div>'
                for l, v in s.items)
            parts.append(f'<section><div class="kpis">{cards}</div></section>')

    meta = []
    if doc.author:
        meta.append(f"<span><b>Author</b> {html.escape(doc.author)}</span>")
    meta.append(f"<span><b>Date</b> {html.escape(doc.date)}</span>")
    meta.append(f"<span><b>Template</b> {html.escape(doc.template)}</span>")
    subtitle = (f'<p class="subtitle">{html.escape(doc.subtitle)}</p>'
                if doc.subtitle else "")
    return (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<meta name='viewport' content='width=device-width,initial-scale=1'>"
        f"<title>{html.escape(doc.title)}</title><style>{_CSS}</style></head><body>"
        f'<div class="page"><header class="doc">'
        f'<h1 class="title">{html.escape(doc.title)}</h1>{subtitle}'
        f'<div class="meta">{"".join(meta)}</div></header>'
        f'{"".join(parts)}'
        f'<footer class="doc">Generated by SciPlotter · {html.escape(doc.date)}</footer>'
        "</div></body></html>"
    )


# --------------------------------------------------------------------- Markdown
def render_markdown(doc: ReportDocument, image_dir=None) -> str:
    """Render Markdown. Figures are written into *image_dir* (if given) and
    linked; otherwise a placeholder is used."""
    import pathlib

    lines = [f"# {doc.title}", ""]
    if doc.subtitle:
        lines += [f"*{doc.subtitle}*", ""]
    meta = []
    if doc.author:
        meta.append(f"**Author:** {doc.author}")
    meta.append(f"**Date:** {doc.date}")
    lines += ["  ".join(meta), ""]
    fig_index = 0
    for s in doc.sections:
        if isinstance(s, TextSection):
            if s.heading:
                lines += [f"## {s.heading}", ""]
            lines += [str(s.body), ""]
        elif isinstance(s, FigureSection):
            if s.title:
                lines += [f"## {s.title}", ""]
            fig_index += 1
            if image_dir and s.png:
                path = pathlib.Path(image_dir) / f"figure_{fig_index}.png"
                path.write_bytes(s.png)
                lines += [f"![{s.caption or s.title}]({path.name})", ""]
            else:
                lines += [f"*(figure: {s.title or 'untitled'})*", ""]
            if s.caption:
                lines += [f"*{s.caption}*", ""]
        elif isinstance(s, TableSection) and s.frame is not None:
            if s.title:
                lines += [f"## {s.title}", ""]
            lines += [_table_markdown(s.frame, s.max_rows), ""]
            if s.caption:
                lines += [f"*{s.caption}*", ""]
        elif isinstance(s, KpiSection) and s.items:
            lines += ["| Metric | Value |", "|---|---|"]
            lines += [f"| {l} | {_fmt(v)} |" for l, v in s.items]
            lines += [""]
    return "\n".join(lines)


def _table_markdown(frame: pd.DataFrame, max_rows: int) -> str:
    shown = frame.head(max_rows)
    header = "| " + " | ".join(str(c) for c in shown.columns) + " |"
    rule = "| " + " | ".join("---" for _ in shown.columns) + " |"
    rows = ["| " + " | ".join(_fmt(v) for v in row) + " |"
            for _, row in shown.iterrows()]
    out = "\n".join([header, rule, *rows])
    if len(frame) > max_rows:
        out += f"\n\n_… {len(frame) - max_rows} more rows ({len(frame)} total)_"
    return out


# -------------------------------------------------------------------------- PDF
def render_pdf(doc: ReportDocument, path: str) -> str:
    """Render the document to a PDF via reportlab platypus. Returns *path*."""
    import io

    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "RTitle", parent=styles["Title"], fontSize=22, spaceAfter=6, textColor=colors.HexColor("#1a2230"))
    h2 = ParagraphStyle(
        "RH2", parent=styles["Heading2"], fontSize=14, spaceBefore=14, spaceAfter=6,
        textColor=colors.HexColor("#1a2230"))
    body = ParagraphStyle("RBody", parent=styles["BodyText"], fontSize=10.5, leading=15)
    muted = ParagraphStyle("RMuted", parent=body, textColor=colors.HexColor("#5a6676"), fontSize=9)

    story = [Paragraph(html.escape(doc.title), title_style)]
    if doc.subtitle:
        story.append(Paragraph(html.escape(doc.subtitle), muted))
    meta = f"Author: {doc.author} &nbsp;·&nbsp; Date: {doc.date}" if doc.author else f"Date: {doc.date}"
    story += [Paragraph(meta, muted), Spacer(1, 0.5 * cm)]

    doc_w = A4[0] - 4 * cm
    for s in doc.sections:
        if isinstance(s, TextSection):
            if s.heading:
                story.append(Paragraph(html.escape(s.heading), h2))
            for para in str(s.body).split("\n\n"):
                if para.strip():
                    story.append(Paragraph(html.escape(para.strip()), body))
        elif isinstance(s, FigureSection) and s.png:
            if s.title:
                story.append(Paragraph(html.escape(s.title), h2))
            img = _scaled_image(io.BytesIO(s.png), doc_w)
            story.append(img)
            if s.caption:
                story.append(Paragraph(html.escape(s.caption), muted))
        elif isinstance(s, TableSection) and s.frame is not None:
            if s.title:
                story.append(Paragraph(html.escape(s.title), h2))
            story.append(_pdf_table(s.frame, s.max_rows, doc_w))
            if s.caption:
                story.append(Paragraph(html.escape(s.caption), muted))
        elif isinstance(s, KpiSection) and s.items:
            data = [[str(l), _fmt(v)] for l, v in s.items]
            t = Table(data, colWidths=[doc_w * 0.6, doc_w * 0.4])
            t.setStyle(TableStyle([
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("TEXTCOLOR", (0, 0), (0, -1), colors.HexColor("#5a6676")),
                ("LINEBELOW", (0, 0), (-1, -1), 0.4, colors.HexColor("#e2e7ee")),
                ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]))
            story.append(t)
        story.append(Spacer(1, 0.3 * cm))

    SimpleDocTemplate(
        str(path), pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm,
        leftMargin=2 * cm, rightMargin=2 * cm,
        title=doc.title, author=doc.author or "SciPlotter",
    ).build(story)
    return str(path)


def _scaled_image(buffer, max_width):
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import Image

    iw, ih = ImageReader(buffer).getSize()
    buffer.seek(0)
    scale = min(1.0, max_width / iw)
    return Image(buffer, width=iw * scale, height=ih * scale)


def render_docx(doc: ReportDocument, path: str) -> str:
    """Render the document to a Word .docx (requires python-docx)."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
    except ImportError as exc:  # pragma: no cover - env without python-docx
        raise RuntimeError(
            "Word export needs python-docx (pip install python-docx)") from exc

    d = Document()
    d.add_heading(doc.title, level=0)
    if doc.subtitle:
        d.add_paragraph(doc.subtitle).italic = True
    meta = f"Author: {doc.author}   Date: {doc.date}" if doc.author else f"Date: {doc.date}"
    m = d.add_paragraph(meta)
    m.runs[0].font.size = Pt(9)
    m.runs[0].font.color.rgb = RGBColor(0x5A, 0x66, 0x76)

    for s in doc.sections:
        if isinstance(s, TextSection):
            if s.heading:
                d.add_heading(s.heading, level=1)
            for para in str(s.body).split("\n\n"):
                if para.strip():
                    d.add_paragraph(para.strip())
        elif isinstance(s, FigureSection) and s.png:
            if s.title:
                d.add_heading(s.title, level=1)
            d.add_picture(io.BytesIO(s.png), width=Inches(6.0))
            if s.caption:
                cap = d.add_paragraph(s.caption)
                cap.runs[0].italic = True
                cap.runs[0].font.size = Pt(9)
        elif isinstance(s, TableSection) and s.frame is not None:
            if s.title:
                d.add_heading(s.title, level=1)
            data = _table_data_docx(s.frame, s.max_rows)
            table = d.add_table(rows=len(data), cols=len(data[0]))
            table.style = "Light Grid Accent 1"
            for r, row in enumerate(data):
                for c, val in enumerate(row):
                    table.rows[r].cells[c].text = val
        elif isinstance(s, KpiSection) and s.items:
            table = d.add_table(rows=len(s.items), cols=2)
            for r, (label, value) in enumerate(s.items):
                table.rows[r].cells[0].text = str(label)
                table.rows[r].cells[1].text = _fmt(value)
    d.save(str(path))
    return str(path)


def _table_data_docx(frame, max_rows):
    shown = frame.head(max_rows)
    data = [[str(c) for c in shown.columns]]
    data += [[_fmt(v) for v in row] for _, row in shown.iterrows()]
    return data


def render_pptx(doc: ReportDocument, path: str) -> str:
    """Render the document to PowerPoint .pptx — a title slide plus one slide
    per figure/table (requires python-pptx)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - env without python-pptx
        raise RuntimeError(
            "PowerPoint export needs python-pptx (pip install python-pptx)") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = doc.title
    if len(slide.placeholders) > 1:
        sub = doc.subtitle or ""
        if doc.author:
            sub = f"{sub}\n{doc.author}".strip()
        slide.placeholders[1].text = f"{sub}\n{doc.date}".strip()

    def _titled(text):
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        return s

    for sec in doc.sections:
        if isinstance(sec, FigureSection) and sec.png:
            s = _titled(sec.title or "Figure")
            s.shapes.add_picture(io.BytesIO(sec.png), Inches(1.2), Inches(1.3),
                                 height=Inches(5.4))
        elif isinstance(sec, TableSection) and sec.frame is not None:
            s = _titled(sec.title or "Table")
            data = _table_data_docx(sec.frame, sec.max_rows)
            rows, cols = len(data), len(data[0])
            tbl = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4),
                                     Inches(11.7), Inches(0.4 * rows)).table
            for r in range(rows):
                for c in range(cols):
                    cell = tbl.cell(r, c)
                    cell.text = data[r][c]
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
        elif isinstance(sec, TextSection) and (sec.heading or sec.body):
            s = _titled(sec.heading or "Notes")
            box = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5))
            box.text_frame.word_wrap = True
            box.text_frame.text = str(sec.body)
            for p in box.text_frame.paragraphs:
                p.font.size = Pt(16)
        elif isinstance(sec, KpiSection) and sec.items:
            s = _titled("Key Results")
            box = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
            tf = box.text_frame
            tf.text = "\n".join(f"{label}:  {_fmt(value)}" for label, value in sec.items)
            for p in tf.paragraphs:
                p.font.size = Pt(20)
    prs.save(str(path))
    return str(path)


def _pdf_table(frame: pd.DataFrame, max_rows: int, width):
    from reportlab.lib import colors
    from reportlab.platypus import Table, TableStyle

    shown = frame.head(max_rows)
    data = [[str(c) for c in shown.columns]]
    data += [[_fmt(v) for v in row] for _, row in shown.iterrows()]
    ncol = max(1, len(shown.columns))
    t = Table(data, colWidths=[width / ncol] * ncol, repeatRows=1)
    t.setStyle(TableStyle([
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f6f8fb")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#1a2230")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("LINEBELOW", (0, 0), (-1, 0), 1, colors.HexColor("#c4cede")),
        ("LINEBELOW", (0, 1), (-1, -1), 0.3, colors.HexColor("#e2e7ee")),
        ("TOPPADDING", (0, 0), (-1, -1), 3), ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    return t
